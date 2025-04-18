# Import dependencies
import argparse
import os
import io
import tempfile
import docx
from tqdm import tqdm
import pptx
import PyPDF2
from xmindparser import xmind_to_dict
import warnings
import msal
import openpyxl
from bs4 import BeautifulSoup
import re
import zipfile
from getpass import getpass
from faster_whisper import WhisperModel
import pandas as pd
import pinecone
from dotenv import load_dotenv
from sharepoint_utils import SharePointFetcher
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_pinecone import PineconeVectorStore
from langchain.vectorstores import Pinecone as LangchainPinecone
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document as LangchainDocument
from pinecone import Pinecone, ServerlessSpec


from typing import List, Optional, Dict, Any
from langchain_community.document_loaders import (
    UnstructuredFileLoader,
    UnstructuredExcelLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredPowerPointLoader,
    CSVLoader
)
from pdf2image import convert_from_bytes
import pytesseract
import warnings

# Class that manages the document loading logic
class SharePointLoader:
    def __init__(self):
        # Initialize text splitter with optimized parameters
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            add_start_index=True,
            separators=["\n\n", "\n", " ", ""] 
        )

        # Initialize Whisper model for audio processing
        self.whisper_model = WhisperModel("base", compute_type="float32")

        # Configure supported file extensions
        self.supported_extensions = {
            "pdf", "docx", "pptx", "xlsx", "xlsm", "csv",
            "mp3", "wav", "m4a", "xmind"
        }

    def load_and_split_file(self, file_name: str, file_data: io.BytesIO) -> List[LangchainDocument]:
        """Load and split a file from SharePoint into LangChain Documents"""
        ext = file_name.split(".")[-1].lower()

        if ext not in self.supported_extensions:
            print(f"[!] Unsupported file type: {file_name}")
            return []

        try:
            # Handle special cases first
            if ext in {"mp3", "wav", "m4a"}:
                return self._process_audio_file(file_data)
            elif ext == "xmind":
                return self._process_xmind_file(file_data)

            # Process other file types
            file_data.seek(0)  # Ensure we're at start of file
            text = self._extract_text(ext, file_data)

            if not text:
                print(f"[!] Empty content extracted from {file_name}")
                return []

            # Create document and split
            doc = LangchainDocument(
                page_content=self._preprocess_text(text),
                metadata={"source": file_name}
            )
            return self.text_splitter.split_documents([doc])

        except Exception as e:
            print(f"[!] Failed to process {file_name}: {e}")
            return []

    def _extract_text(self, ext: str, file_data: io.BytesIO) -> str:
        """Extract and format text content based on file extension."""
        file_data.seek(0)

        if ext == "pdf":
            reader = PyPDF2.PdfReader(file_data)
            extracted = "\n".join(page.extract_text() or "" for page in reader.pages)

            if not extracted.strip():  # Use OCR to extract text from image (if any)
                file_data.seek(0)
                images = convert_from_bytes(file_data.read())
                ocr_text = "\n".join(pytesseract.image_to_string(img) for img in images)
                return ocr_text

            return extracted

        elif ext == "docx":
            doc = docx.Document(file_data)
            return "\n".join(para.text for para in doc.paragraphs)

        elif ext == "pptx":
            prs = pptx.Presentation(file_data)
            return "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            )

        elif ext in {"xlsx", "xlsm", "csv"}:
            return self._extract_tabular_text(file_data, ext)

        return ""

    def _extract_tabular_text(self, file_data: io.BytesIO, ext: str) -> str:
        try:
            if ext == "csv":
                df = pd.read_csv(file_data)
                if df.empty or df.shape[1] == 0:
                    return ""
                return "\n".join(self._format_dataframe_rows(df))

            else:
                with warnings.catch_warnings():
                    warnings.simplefilter("ignore", UserWarning)
                    df_dict = pd.read_excel(file_data, sheet_name=None, engine="openpyxl")

                if not isinstance(df_dict, dict):
                    return ""

                all_text = []
                for sheet_name, sheet_df in df_dict.items():
                    if sheet_df.empty or sheet_df.shape[1] == 0:
                        continue
                    all_text.append(f"Sheet: {sheet_name}")
                    all_text.extend(self._format_dataframe_rows(sheet_df))

                return "\n\n".join(all_text)

        except Exception as e:
            print(f"[!] Failed to read tabular file: {e}")
            return ""


    def _format_dataframe_rows(self, df: pd.DataFrame) -> List[str]:
        """Convert each row in a DataFrame to a structured string."""
        df.dropna(how='all', inplace=True)
        df.dropna(axis=1, how='all', inplace=True)
        df = df.astype(str).replace("nan", "N/A").replace("NaT", "N/A")

        formatted_rows = []
        for _, row in df.iterrows():
            row_text = "\n".join(f"{str(col).strip()}: {str(val).strip() if isinstance(val, str) else str(val)}" for col, val in row.items())
            formatted_rows.append(row_text)

        return formatted_rows


    def _process_audio_file(self, file_data: io.BytesIO) -> List[LangchainDocument]:
        """Process audio files with Whisper"""
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as tmp:
            file_data.seek(0)
            tmp.write(file_data.read())
            tmp_path = tmp.name

        try:
            segments, _ = self.whisper_model.transcribe(tmp_path)
            full_text = " ".join(segment.text for segment in segments)
            doc = LangchainDocument(
                page_content=self._preprocess_text(full_text),
                metadata={"source": "audio_transcription"}
            )
            return self.text_splitter.split_documents([doc])
        finally:
            os.unlink(tmp_path)

    def _process_xmind_file(self, file_data: io.BytesIO) -> List[LangchainDocument]:
        """Process XMind files with custom parser"""
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xmind") as tmp:
            file_data.seek(0)
            tmp.write(file_data.read())
            tmp_path = tmp.name

        try:
            content = xmind_to_dict(tmp_path)
            text_chunks = []

            def traverse(topics, prefix=""):
                for topic in topics:
                    title = topic.get('title', '')
                    if title:  # Only add non-empty titles
                        text_chunks.append(prefix + title)
                    if 'topics' in topic:
                        traverse(topic['topics'], prefix + "  ")

            for sheet in content:
                traverse(sheet.get('topic', {}).get('topics', []))

            full_text = "\n".join(text_chunks)
            doc = LangchainDocument(
                page_content=self._preprocess_text(full_text),
                metadata={"source": "xmind_file"}
            )
            return self.text_splitter.split_documents([doc])
        except Exception as e:
            print(f"[!] Failed to parse XMind file: {e}")
            return []
        finally:
            os.unlink(tmp_path)

    def _preprocess_text(self, text: str) -> str:
        """Clean and normalize text content for better ingestion"""
        # Lowercase everything
        text = text.lower()

        # Replace common spreadsheet artifacts with space
        text = re.sub(r'\b(nan|n/a|na|none|null)\b', ' ', text)

        # Remove time-like values (e.g., 00:00:00)
        text = re.sub(r'\b\d{1,2}:\d{2}:\d{2}\b', ' ', text)

        # Remove standalone numeric entries or sequences (e.g., 0.00, 123.0)
        text = re.sub(r'\b\d+(\.\d+)?\b', ' ', text)

        # Remove excessive booleans
        text = re.sub(r'\b(true|false|yes|no)\b', ' ', text)

        # Remove repeated phrases (simple deduplication)
        words = text.split()
        deduped = []
        last_word = None
        for word in words:
            if word != last_word:
                deduped.append(word)
            last_word = word
        text = " ".join(deduped)

        # Replace newline indicators or multiple spaces
        text = text.replace("\\n", " ")
        text = re.sub(r'\s+', ' ', text).strip()

        return text


# Function to ingest files to pinecone
def ingest_files_to_pinecone(
    files: Dict[str, io.BytesIO],
    loader: SharePointLoader,
    vectorstore: LangchainPinecone,
    batch_size: int = 100
) -> None:
    """
    Complete ingestion pipeline with proper error handling and batching

    Args:
        files: Dictionary of {file_path: file_data} from SharePoint
        loader: Initialized SharePointLoader instance
        vectorstore: Initialized LangchainPinecone vector store
        batch_size: Number of documents to process in each batch
    """
    all_documents = []
    failed_files = []

    # First pass: Process all files and collect documents
    for file_path, file_data in tqdm(files.items(), desc="Processing files"):
        try:
            docs = loader.load_and_split_file(file_path, file_data)
            if not docs:
                print(f"\n[!] No content extracted from {file_path}")
                continue

            # Add source metadata to all chunks
            for doc in docs:
                doc.metadata.update({
                    "source": file_path,
                    "file_type": file_path.split(".")[-1].lower()
                })

            all_documents.extend(docs)

        except Exception as e:
            print(f"\n[!] Failed to process {file_path}: {str(e)}")
            failed_files.append(file_path)

    if not all_documents:
        print("\n[!] No documents to ingest")
        return

    # Second pass: Batch ingest to Pinecone
    print(f"\nStarting ingestion of {len(all_documents)} chunks to Pinecone")

    successful_chunks = 0
    for i in tqdm(range(0, len(all_documents), batch_size),
                desc="Ingesting to Pinecone"):
        batch = all_documents[i:i + batch_size]
        try:
            vectorstore.add_documents(batch)
            successful_chunks += len(batch)
        except Exception as e:
            print(f"\n[!] Failed to ingest batch {i//batch_size + 1}: {str(e)}")

    # Print summary
    print("\nIngestion Summary:")
    print(f"- Successfully processed {len(files) - len(failed_files)}/{len(files)} files")
    print(f"- Successfully ingested {successful_chunks}/{len(all_documents)} chunks")

    if failed_files:
        print("\nFiles that failed processing:")
        for file_path in failed_files[:5]:  # Show first 5 to avoid flooding
            print(f"  - {file_path}")
        if len(failed_files) > 5:
            print(f"  ... and {len(failed_files) - 5} more")
